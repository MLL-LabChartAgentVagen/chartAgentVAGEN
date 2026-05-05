"""Generate ChartAgent proposal slides as editable .pptx.

All diagrams are built from native python-pptx shapes, so every element
remains editable in PowerPoint. No embedded images, no rasterization.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ─────────────────────────── Design tokens ───────────────────────────

NAVY       = RGBColor(0x0F, 0x3A, 0x5F)   # primary
NAVY_DARK  = RGBColor(0x08, 0x25, 0x3E)
ACCENT     = RGBColor(0xD9, 0x72, 0x1E)   # highlight orange
TEAL       = RGBColor(0x2C, 0x8C, 0x8A)   # secondary
TEXT_DARK  = RGBColor(0x1F, 0x25, 0x37)
TEXT_MID   = RGBColor(0x55, 0x5F, 0x6E)
TEXT_MUTED = RGBColor(0x8A, 0x94, 0xA6)
BORDER     = RGBColor(0xD8, 0xDE, 0xE8)
SOFT_BG    = RGBColor(0xF6, 0xF8, 0xFB)
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)
CHIP_BG    = RGBColor(0xE8, 0xEF, 0xF7)
GRAY_BG    = RGBColor(0xF0, 0xF2, 0xF5)
CODE_BG    = RGBColor(0x11, 0x1B, 0x2E)
CODE_FG    = RGBColor(0xE6, 0xED, 0xF6)
CODE_KEY   = RGBColor(0xFF, 0xB3, 0x6B)   # keywords
CODE_STR   = RGBColor(0x9F, 0xE5, 0xA0)   # strings
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─────────────────────────── Helpers ───────────────────────────

def _set_fill(shape, color, transparent=False):
    if transparent:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _set_line(shape, color=None, width_pt=0.75, invisible=False):
    if invisible:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color if color else BORDER
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=BORDER, line_w=0.75,
             rounded=False, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        shp.adjustments[0] = 0.08
    _set_fill(shp, fill, transparent=(fill is None))
    if line is None:
        _set_line(shp, invisible=True)
    else:
        _set_line(shp, line, line_w)
    if not shadow:
        shp.shadow.inherit = False
        # disable default shadow via XML
        sppr = shp.fill._xPr
        for eff in sppr.findall(qn('a:effectLst')):
            sppr.remove(eff)
        eff_lst = etree.SubElement(sppr, qn('a:effectLst'))
    return shp


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, italic=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs = list of (text, dict(size, bold, italic, color, font))."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", Emu(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    first = True
    for text, style in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            first = True
            continue
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = style.get("font", FONT)
        f.size = Pt(style.get("size", 12))
        f.bold = style.get("bold", False)
        f.italic = style.get("italic", False)
        f.color.rgb = style.get("color", TEXT_DARK)
        first = False
    return tb


def add_arrow_right(slide, x, y, w, h, *, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    # slimmer arrowhead: adj1 = head width, adj2 = head length
    shp.adjustments[0] = 0.50
    shp.adjustments[1] = 0.30
    _set_fill(shp, fill)
    _set_line(shp, invisible=True) if line is None else _set_line(shp, line, 0.75)
    return shp


def add_line(slide, x1, y1, x2, y2, *, color=TEXT_MID, width_pt=1.25,
             dashed=False, arrow_end=False, arrow_start=False):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    ln = connector.line._get_or_add_ln()
    if dashed:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    if arrow_end:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    if arrow_start:
        head = etree.SubElement(ln, qn('a:headEnd'))
        head.set('type', 'triangle')
        head.set('w', 'med')
        head.set('len', 'med')
    return connector


def add_chip(slide, x, y, text, *, size=10, fill=CHIP_BG, color=NAVY,
             pad_x=0.12, height=0.32, bold=False):
    w = Inches(max(0.7, 0.10 * len(text) + 0.30))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(height))
    shp.adjustments[0] = 0.5
    _set_fill(shp, fill)
    _set_line(shp, invisible=True)
    tf = shp.text_frame
    tf.margin_left = Inches(pad_x); tf.margin_right = Inches(pad_x)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shp, w


def fill_shape_text(shape, text, *, size=12, bold=False, color=TEXT_DARK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    font=FONT, italic=False):
    tf = shape.text_frame
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title, subtitle=None, accent=True):
    # page title
    add_text(slide, Inches(0.55), Inches(0.35), Inches(11.5), Inches(0.55),
             title, size=26, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.90), Inches(11.5), Inches(0.35),
                 subtitle, size=13, color=TEXT_MID, italic=True)
    # accent underline
    if accent:
        bar = add_rect(slide, Inches(0.55), Inches(1.25), Inches(0.5), Inches(0.05),
                       fill=ACCENT, line=None)


def add_footer(slide, page, total=15, tag="ChartAgent · Atomic-Grain Programmatic Data Synthesis"):
    add_text(slide, Inches(0.55), Inches(7.10), Inches(10.0), Inches(0.3),
             tag, size=9, color=TEXT_MUTED)
    add_text(slide, Inches(11.5), Inches(7.10), Inches(1.3), Inches(0.3),
             f"{page} / {total}", size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_bg(slide, WHITE)
    return slide


# ─────────────────────────── Slide 1 — Title ───────────────────────────

def slide_title(prs):
    s = blank_slide(prs)

    # left accent column
    add_rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, fill=NAVY, line=None)
    add_rect(s, Inches(0.4), Inches(0), Inches(0.08), SLIDE_H, fill=ACCENT, line=None)

    # eyebrow
    add_text(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(0.3),
             "RESEARCH PROPOSAL", size=12, bold=True, color=ACCENT)

    # main title
    add_text(s, Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.3),
             "ChartAgent", size=62, bold=True, color=NAVY)

    add_text(s, Inches(1.0), Inches(3.25), Inches(11.5), Inches(0.9),
             "Atomic-Grain Programmatic Data Synthesis\nfor Chart Understanding Benchmarks",
             size=26, color=TEXT_DARK, line_spacing=1.2)

    # underline
    add_rect(s, Inches(1.0), Inches(4.55), Inches(1.6), Inches(0.04),
             fill=ACCENT, line=None)

    # three pillar chips
    pillars = ["Operator Algebra", "Agentic Data Simulator", "Table Amortization"]
    x = Inches(1.0); y = Inches(4.9)
    for i, p in enumerate(pillars):
        shp, w = add_chip(s, x, y, p, size=13, fill=NAVY, color=WHITE,
                          height=0.46, pad_x=0.22, bold=True)
        x = x + w + Inches(0.20)

    # tagline
    add_text(s, Inches(1.0), Inches(5.85), Inches(11.5), Inches(0.45),
             "One Master Table  →  10–30+ Coherent Multi-Chart Tasks",
             size=16, color=TEXT_MID, italic=True)

    # corner meta
    add_text(s, Inches(1.0), Inches(6.6), Inches(11.5), Inches(0.35),
             "Deterministic · Type-Safe · Cross-Chart Consistent by Construction",
             size=11, color=TEXT_MUTED)


# ─────────────────────────── Slide 2 — Motivation ───────────────────────────

def slide_motivation(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Why existing chart QA benchmarks fall short",
                  "Three structural deficits shared by ChartQA, PlotQA, ChartBench")

    card_w, card_h = Inches(3.95), Inches(5.0)
    gap = Inches(0.20)
    y = Inches(1.65)
    xs = [Inches(0.55), Inches(0.55) + card_w + gap, Inches(0.55) + 2*(card_w + gap)]

    cards = [
        {
            "num": "01",
            "title": "Shallow data depth",
            "body": "Public datasets store only the 5-row aggregate shown in the chart.",
            "visual": "agg_table",
            "pain": "No drill-down, no provenance, no long-tail statistics.",
        },
        {
            "num": "02",
            "title": "Arithmetic hallucinations",
            "body": "LLM-written JSON introduces numbers that don't add up.",
            "visual": "broken_math",
            "pain": "Validation catches syntax — not semantic impossibility.",
        },
        {
            "num": "03",
            "title": "No cross-chart consistency",
            "body": "A pie and a bar from the same domain disagree on totals.",
            "visual": "mismatch",
            "pain": "Multi-chart reasoning becomes impossible to evaluate.",
        },
    ]

    for x, c in zip(xs, cards):
        # card
        add_rect(s, x, y, card_w, card_h, fill=CARD_BG, line=BORDER,
                 line_w=0.75, rounded=True)
        # big number
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(1.2), Inches(0.7),
                 c["num"], size=36, bold=True, color=ACCENT)
        # title
        add_text(s, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.5),
                 c["title"], size=18, bold=True, color=NAVY)
        # body
        add_text(s, x + Inches(0.25), y + Inches(1.45), card_w - Inches(0.5), Inches(0.9),
                 c["body"], size=12, color=TEXT_DARK, line_spacing=1.3)

        # visual zone
        vx = x + Inches(0.35); vy = y + Inches(2.4)
        vw = card_w - Inches(0.7); vh = Inches(1.7)
        if c["visual"] == "agg_table":
            # mini table: 3 rows × 2 cols
            tw, th = vw, Inches(1.35)
            tx, ty = vx, vy + Inches(0.15)
            headers = ["region", "revenue"]
            rows = [("North", "$4.2M"), ("South", "$3.1M"), ("East", "$2.8M")]
            h_each = th / 4
            # header row
            add_rect(s, tx, ty, tw, h_each, fill=NAVY, line=None)
            for i, htx in enumerate(headers):
                add_text(s, tx + i*(tw/2) + Inches(0.12), ty + Inches(0.04),
                         tw/2 - Inches(0.2), h_each, htx,
                         size=10, bold=True, color=WHITE)
            for r_i, (a, b) in enumerate(rows, start=1):
                add_rect(s, tx, ty + h_each*r_i, tw, h_each,
                         fill=WHITE if r_i % 2 else SOFT_BG, line=BORDER, line_w=0.5)
                add_text(s, tx + Inches(0.12), ty + h_each*r_i + Inches(0.04),
                         tw/2 - Inches(0.2), h_each, a, size=10, color=TEXT_DARK)
                add_text(s, tx + tw/2 + Inches(0.12), ty + h_each*r_i + Inches(0.04),
                         tw/2 - Inches(0.2), h_each, b, size=10, color=TEXT_DARK)
            add_text(s, vx, vy + th + Inches(0.25), vw, Inches(0.3),
                     "only aggregates — no atomic events",
                     size=9, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        elif c["visual"] == "broken_math":
            # show "10 + 15 ≠ 30" with crossed-out result
            add_runs(s, vx, vy + Inches(0.35), vw, Inches(0.7),
                     [("10  +  15  =  ", {"size": 26, "bold": True, "color": TEXT_DARK, "font": MONO}),
                      ("30", {"size": 26, "bold": True, "color": ACCENT, "font": MONO})],
                     align=PP_ALIGN.CENTER)
            # red cross line
            add_line(s, vx + vw/2 + Inches(0.6), vy + Inches(0.55),
                     vx + vw/2 + Inches(1.2), vy + Inches(0.95),
                     color=ACCENT, width_pt=2.0)
            add_text(s, vx, vy + Inches(1.25), vw, Inches(0.35),
                     "LLM-produced numbers silently break arithmetic",
                     size=9, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        else:  # mismatch
            # two mini charts with mismatched totals
            b1x = vx + Inches(0.15); b1y = vy + Inches(0.1)
            add_text(s, b1x, b1y, Inches(1.1), Inches(0.3),
                     "Pie  →  $10.0M", size=10, color=TEXT_DARK, bold=True)
            # circle
            pie = s.shapes.add_shape(MSO_SHAPE.OVAL, b1x + Inches(0.15),
                                      b1y + Inches(0.35), Inches(0.8), Inches(0.8))
            _set_fill(pie, NAVY); _set_line(pie, invisible=True)
            b2x = vx + Inches(1.95); b2y = vy + Inches(0.1)
            add_text(s, b2x, b2y, Inches(1.1), Inches(0.3),
                     "Bar  →  $11.4M", size=10, color=TEXT_DARK, bold=True)
            # bars
            for i, hgt in enumerate([0.35, 0.55, 0.70]):
                add_rect(s, b2x + Inches(0.1 + i*0.28), b2y + Inches(1.15 - hgt),
                         Inches(0.22), Inches(hgt), fill=TEAL, line=None)
            # mismatch label
            add_text(s, vx, vy + Inches(1.3), vw, Inches(0.3),
                     "same data  ·  totals disagree", size=9, italic=True,
                     color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

        # bottom pain quote
        add_rect(s, x + Inches(0.25), y + card_h - Inches(0.85),
                 card_w - Inches(0.5), Inches(0.65),
                 fill=SOFT_BG, line=None, rounded=True)
        add_text(s, x + Inches(0.35), y + card_h - Inches(0.75),
                 card_w - Inches(0.7), Inches(0.5),
                 c["pain"], size=10, italic=True, color=TEXT_MID, line_spacing=1.25)

    add_footer(s, 2)


# ─────────────────────────── Slide 3 — Prior work ───────────────────────────

def slide_prior_work(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Comparison with prior work",
                  "Every prior benchmark fails ≥ 4 of the 7 axes we care about")

    headers = ["Dimension", "ChartQA", "PlotQA", "ChartBench", "Ours"]
    rows = [
        ["Data source",             "Web-crawled",   "Templates",       "Hybrid",         "LLM-authored SDK programs"],
        ["Data depth",              "Aggregates",    "Aggregates",      "Aggregates",     "Atomic-grain fact tables"],
        ["Distribution control",    "None",          "Uniform/Normal",  "Semi-fixed",     "Mixtures / copulas / conditional"],
        ["Cross-chart consistency", "✗",             "✗",               "✗",              "✓  (shared Master Table)"],
        ["Error correction",        "Manual",        "None",            "Heuristic",      "Code-execution feedback loop"],
        ["Multi-chart reasoning",   "None",          "None",            "Limited",        "Dashboard-level cross-chart QA"],
        ["Reproducibility",         "✗",             "✓",               "Partial",        "✓  (SDK + seed)"],
    ]

    col_xs = [Inches(0.55), Inches(3.8), Inches(5.6), Inches(7.4), Inches(9.4)]
    col_ws = [Inches(3.20), Inches(1.75), Inches(1.75), Inches(1.95), Inches(3.35)]
    row_h = Inches(0.60)
    y = Inches(1.75)

    # Header
    for i, (cx, cw, ht) in enumerate(zip(col_xs, col_ws, headers)):
        fill = NAVY_DARK if i == 4 else NAVY
        add_rect(s, cx, y, cw, row_h, fill=fill, line=None)
        add_text(s, cx + Inches(0.12), y, cw - Inches(0.2), row_h,
                 ht, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)

    for r_i, row in enumerate(rows):
        ry = y + row_h * (r_i + 1)
        for i, (cx, cw, val) in enumerate(zip(col_xs, col_ws, row)):
            is_ours = (i == 4)
            fill = CHIP_BG if is_ours else (SOFT_BG if r_i % 2 == 0 else WHITE)
            add_rect(s, cx, ry, cw, row_h, fill=fill, line=BORDER, line_w=0.5)
            color = NAVY if is_ours else TEXT_DARK
            bold = True if is_ours or i == 0 else False
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
            size = 11 if i > 0 else 11.5
            add_text(s, cx + Inches(0.12), ry, cw - Inches(0.2), row_h,
                     val, size=size, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE, align=align)

    # caption
    add_text(s, Inches(0.55), Inches(6.55), Inches(12), Inches(0.4),
             "Key difference: we control data depth, distribution, and cross-chart semantics — not just the surface image.",
             size=12, italic=True, color=TEXT_MID)

    add_footer(s, 3)


# ─────────────────────────── Slide 4 — Three contributions ───────────────────────────

def slide_contributions(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Three interlocking contributions",
                  "A formal substrate → a programmable simulator → a deterministic amortizer")

    card_w, card_h = Inches(4.0), Inches(4.9)
    gap = Inches(0.15)
    y = Inches(1.75)
    xs = [Inches(0.55), Inches(0.55) + card_w + gap, Inches(0.55) + 2*(card_w + gap)]

    items = [
        {
            "num": "1",
            "tag": "FORMAL",
            "title": "Operator Algebra",
            "body": "Every question is a typed pipeline over Set / Scalar / Bridge operators.",
            "bullets": ["16 operators, 2 types (V, S)",
                        "Difficulty = # ops",
                        "Composable multi-hop reasoning"],
        },
        {
            "num": "2",
            "tag": "EXECUTABLE",
            "title": "Agentic Data Simulator",
            "body": "LLM writes Python calling a type-safe SDK — not JSON, not raw numbers.",
            "bullets": ["Closed-form DGP per measure",
                        "Code-execution feedback loop",
                        "3-layer validator, no LLM re-call"],
        },
        {
            "num": "3",
            "tag": "AMORTIZED",
            "title": "Table Amortization",
            "body": "One Master Table → 10–30+ charts with arithmetic consistency guaranteed.",
            "bullets": ["16 chart types × 7 relationships",
                        "Deterministic SQL projection",
                        "Cross-chart QA by construction"],
        },
    ]
    accent_colors = [TEAL, NAVY, ACCENT]

    for x, c, ac in zip(xs, items, accent_colors):
        add_rect(s, x, y, card_w, card_h, fill=CARD_BG, line=BORDER,
                 line_w=0.75, rounded=True)
        # top accent bar
        add_rect(s, x, y, card_w, Inches(0.12), fill=ac, line=None, rounded=False)

        # giant number
        add_text(s, x + Inches(0.35), y + Inches(0.35), Inches(0.9), Inches(1.4),
                 c["num"], size=66, bold=True, color=ac)
        # tag
        add_text(s, x + Inches(1.30), y + Inches(0.55), Inches(2.5), Inches(0.3),
                 c["tag"], size=10, bold=True, color=TEXT_MUTED)
        # title (on one line now)
        add_text(s, x + Inches(1.30), y + Inches(0.85), card_w - Inches(1.55), Inches(0.9),
                 c["title"], size=19, bold=True, color=NAVY, line_spacing=1.1)
        # body
        add_text(s, x + Inches(0.35), y + Inches(2.0), card_w - Inches(0.7), Inches(1.0),
                 c["body"], size=13, color=TEXT_DARK, line_spacing=1.35)

        # divider
        add_line(s, x + Inches(0.35), y + Inches(3.10),
                 x + card_w - Inches(0.35), y + Inches(3.10),
                 color=BORDER, width_pt=0.75)

        # bullets
        by = y + Inches(3.30)
        for bi, b in enumerate(c["bullets"]):
            # bullet dot
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.4), by + Inches(0.12),
                                     Inches(0.09), Inches(0.09))
            _set_fill(dot, ac); _set_line(dot, invisible=True)
            add_text(s, x + Inches(0.60), by + Inches(0.01),
                     card_w - Inches(0.9), Inches(0.35),
                     b, size=12, color=TEXT_DARK)
            by = by + Inches(0.45)

    add_footer(s, 4)


# ─────────────────────────── Slide 5 — 4-Phase Pipeline ───────────────────────────

def slide_pipeline(prs):
    s = blank_slide(prs)
    add_title_bar(s, "The 4-phase pipeline",
                  "LLM called only in phases 0 – 2.  Phase 3 is fully deterministic.")

    # Four stage boxes
    stage_y = Inches(2.2)
    stage_h = Inches(2.3)
    sw = Inches(2.7)
    gap = Inches(0.35)
    xs = [Inches(0.55) + i*(sw + gap) for i in range(4)]

    stages = [
        ("PHASE 0", "Domain Pool\nConstruction",
         "200+ fine-grained\ndomains, cached",
         "LLM · one-time", NAVY),
        ("PHASE 1", "Scenario\nContextualization",
         "Entities, metrics,\ntemporal grain",
         "LLM · per-sample", NAVY),
        ("PHASE 2", "Agentic Data\nSimulator (SDK)",
         "Master Fact Table\n+ Schema Metadata",
         "LLM · per-sample", ACCENT),
        ("PHASE 3", "View Amortization\n+ QA Generation",
         "10–30+ coherent\nmulti-chart tasks",
         "Deterministic", TEAL),
    ]

    for x, (tag, title, out, llm, col) in zip(xs, stages):
        # main card
        add_rect(s, x, stage_y, sw, stage_h, fill=CARD_BG, line=BORDER,
                 line_w=0.75, rounded=True)
        # top color bar
        add_rect(s, x, stage_y, sw, Inches(0.45), fill=col, line=None)
        add_text(s, x + Inches(0.2), stage_y + Inches(0.08),
                 sw - Inches(0.4), Inches(0.3),
                 tag, size=11, bold=True, color=WHITE)
        # title
        add_text(s, x + Inches(0.2), stage_y + Inches(0.6),
                 sw - Inches(0.4), Inches(0.75),
                 title, size=15, bold=True, color=NAVY, line_spacing=1.15)
        # divider
        add_line(s, x + Inches(0.3), stage_y + Inches(1.35),
                 x + sw - Inches(0.3), stage_y + Inches(1.35),
                 color=BORDER)
        # output
        add_text(s, x + Inches(0.2), stage_y + Inches(1.40),
                 sw - Inches(0.4), Inches(0.65),
                 out, size=11, color=TEXT_DARK, line_spacing=1.25, italic=True)
        # llm badge
        badge_y = stage_y + stage_h - Inches(0.45)
        llm_fill = GRAY_BG if "Deterministic" in llm else CHIP_BG
        llm_color = TEAL if "Deterministic" in llm else NAVY
        add_rect(s, x + Inches(0.2), badge_y,
                 sw - Inches(0.4), Inches(0.3),
                 fill=llm_fill, line=None, rounded=True)
        add_text(s, x + Inches(0.2), badge_y, sw - Inches(0.4), Inches(0.3),
                 llm, size=10, bold=True, color=llm_color,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # Connecting arrows
    for i in range(3):
        ax = xs[i] + sw + Inches(0.02)
        ay = stage_y + stage_h / 2 - Inches(0.15)
        aw = gap - Inches(0.04)
        add_arrow_right(s, ax, ay, aw, Inches(0.3), fill=NAVY)

    # Feedback loop on phase 2 (curve shown as two segments)
    p2_x = xs[2]; p2_w = sw
    loop_y = Inches(4.85)
    add_line(s, p2_x + p2_w - Inches(0.7), stage_y + stage_h,
             p2_x + p2_w - Inches(0.7), loop_y,
             color=ACCENT, width_pt=1.5, dashed=True)
    add_line(s, p2_x + p2_w - Inches(0.7), loop_y,
             p2_x + Inches(0.5), loop_y,
             color=ACCENT, width_pt=1.5, dashed=True)
    add_line(s, p2_x + Inches(0.5), loop_y,
             p2_x + Inches(0.5), stage_y + stage_h,
             color=ACCENT, width_pt=1.5, dashed=True, arrow_end=True)
    add_text(s, p2_x + Inches(0.0), loop_y + Inches(0.04),
             p2_w, Inches(0.35),
             "Code-execution feedback  (max 3 retries)",
             size=10, italic=True, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)

    # Output band at the bottom
    out_y = Inches(5.55)
    add_rect(s, Inches(0.55), out_y, Inches(12.22), Inches(1.0),
             fill=SOFT_BG, line=BORDER, line_w=0.5, rounded=True)
    add_text(s, Inches(0.8), out_y + Inches(0.10),
             Inches(12), Inches(0.35),
             "OUTPUT", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), out_y + Inches(0.38),
             Inches(12), Inches(0.55),
             "{ Chart Images, Questions, Answers, Reasoning Chains }  ×  N",
             size=16, bold=True, color=NAVY, font=MONO)

    add_footer(s, 5)


# ─────────────────────────── Slide 6 — Phase 0 + 1 ───────────────────────────

def slide_phase_01(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Phase 0 & 1:  From domain pool to scenario",
                  "Chart types are never mentioned — data is born from business needs.")

    # two columns
    col_w = Inches(6.0)
    left_x = Inches(0.55); right_x = Inches(6.78)
    y = Inches(1.75)
    col_h = Inches(5.1)

    # LEFT — Phase 0
    add_rect(s, left_x, y, col_w, col_h, fill=CARD_BG, line=BORDER,
             line_w=0.75, rounded=True)
    add_rect(s, left_x, y, col_w, Inches(0.55), fill=NAVY, line=None)
    add_text(s, left_x + Inches(0.25), y + Inches(0.12),
             col_w - Inches(0.5), Inches(0.35),
             "PHASE 0  ·  Domain Pool  (one-time, cached)",
             size=14, bold=True, color=WHITE)

    # topic chips
    ty = y + Inches(0.8)
    add_text(s, left_x + Inches(0.25), ty, col_w - Inches(0.5), Inches(0.3),
             "15+ super-sectors", size=11, bold=True, color=TEXT_MID)
    topics = ["Healthcare", "Finance", "Retail", "Transportation",
              "Energy", "Manufacturing", "Education", "Agriculture"]
    cx = left_x + Inches(0.25); cy = ty + Inches(0.4)
    for t in topics:
        shp, cw = add_chip(s, cx, cy, t, size=10, fill=CHIP_BG, color=NAVY, height=0.32)
        cx = cx + cw + Inches(0.08)
        if cx > left_x + col_w - Inches(1.5):
            cx = left_x + Inches(0.25); cy = cy + Inches(0.42)

    # stats row
    stat_y = y + Inches(2.35)
    add_line(s, left_x + Inches(0.25), stat_y - Inches(0.15),
             left_x + col_w - Inches(0.25), stat_y - Inches(0.15),
             color=BORDER)
    stat_items = [("213", "domains"), ("15+", "topics"),
                  ("3", "complexity tiers"), ("0.82", "diversity score")]
    sw = (col_w - Inches(0.5)) / 4
    for i, (num, lab) in enumerate(stat_items):
        sx = left_x + Inches(0.25) + sw * i
        add_text(s, sx, stat_y, sw, Inches(0.45),
                 num, size=26, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER)
        add_text(s, sx, stat_y + Inches(0.50), sw, Inches(0.3),
                 lab, size=10, color=TEXT_MID, align=PP_ALIGN.CENTER)

    # complexity bar
    bar_y = y + Inches(3.5)
    add_text(s, left_x + Inches(0.25), bar_y, col_w - Inches(0.5), Inches(0.3),
             "Complexity-balanced stratification",
             size=11, bold=True, color=TEXT_MID)
    seg_y = bar_y + Inches(0.4)
    seg_h = Inches(0.35)
    parts = [("simple", 71, NAVY), ("medium", 72, TEAL), ("complex", 70, ACCENT)]
    total = sum(p[1] for p in parts)
    sx = left_x + Inches(0.25)
    max_w = col_w - Inches(0.5)
    for name, cnt, c in parts:
        pw = max_w * (cnt / total)
        add_rect(s, sx, seg_y, pw, seg_h, fill=c, line=None)
        add_text(s, sx, seg_y, pw, seg_h,
                 f"{name} · {cnt}", size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        sx = sx + pw

    # RIGHT — Phase 1
    add_rect(s, right_x, y, col_w, col_h, fill=CARD_BG, line=BORDER,
             line_w=0.75, rounded=True)
    add_rect(s, right_x, y, col_w, Inches(0.55), fill=TEAL, line=None)
    add_text(s, right_x + Inches(0.25), y + Inches(0.12),
             col_w - Inches(0.5), Inches(0.35),
             "PHASE 1  ·  Scenario Contextualization",
             size=14, bold=True, color=WHITE)

    # Sample domain → scenario flow
    add_text(s, right_x + Inches(0.25), y + Inches(0.8),
             col_w - Inches(0.5), Inches(0.3),
             "Sampled domain  →  concrete realistic scenario",
             size=11, bold=True, color=TEXT_MID)

    # JSON preview card (code block)
    jy = y + Inches(1.2)
    jh = Inches(3.70)
    add_rect(s, right_x + Inches(0.25), jy, col_w - Inches(0.5), jh,
             fill=CODE_BG, line=None, rounded=True)
    json_lines = [
        ('{', CODE_FG),
        ('  "scenario_title":', CODE_KEY),
        ('    "2024 H1 Shanghai Metro Ridership Log",', CODE_STR),
        ('  "data_context":', CODE_KEY),
        ('    "Shanghai Transport Commission collected ... ",', CODE_STR),
        ('  "temporal_granularity":', CODE_KEY),
        ('    "daily",', CODE_STR),
        ('  "key_entities": [', CODE_KEY),
        ('     "Line 1", "Line 2", "Line 8", ...', CODE_STR),
        ('  ],', CODE_FG),
        ('  "key_metrics": [', CODE_KEY),
        ('     { "name": "daily_ridership",', CODE_STR),
        ('       "unit": "10k passengers",', CODE_STR),
        ('       "range": [5, 120] },', CODE_STR),
        ('     { "name": "on_time_rate", ... }', CODE_STR),
        ('  ],', CODE_FG),
        ('  "target_rows": 900', CODE_KEY),
        ('}', CODE_FG),
    ]
    ly = jy + Inches(0.15)
    for text, col in json_lines:
        add_text(s, right_x + Inches(0.45), ly, col_w - Inches(0.7), Inches(0.22),
                 text, size=10.5, color=col, font=MONO, line_spacing=1.0)
        ly = ly + Inches(0.195)

    # caption
    add_text(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.3),
             "Chart-type isolation is the key to breaking “template disease” in existing benchmarks.",
             size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_footer(s, 6)


# ─────────────────────────── Slide 7 — Paradigm shift ───────────────────────────

def slide_paradigm(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Phase 2 paradigm shift:  code-as-DGP",
                  "From brittle JSON configs to executable, type-safe Python.")

    # Two columns
    left_x, right_x = Inches(0.55), Inches(6.92)
    col_w = Inches(5.85)
    y = Inches(1.75)
    h = Inches(5.1)

    # LEFT — old way (faded)
    add_rect(s, left_x, y, col_w, h, fill=SOFT_BG, line=BORDER,
             line_w=0.75, rounded=True)
    add_text(s, left_x + Inches(0.3), y + Inches(0.3),
             col_w - Inches(0.6), Inches(0.35),
             "BEFORE", size=11, bold=True, color=TEXT_MUTED)
    add_text(s, left_x + Inches(0.3), y + Inches(0.6),
             col_w - Inches(0.6), Inches(0.6),
             "LLM-as-Data-Generator", size=20, bold=True, color=TEXT_MID)
    add_text(s, left_x + Inches(0.3), y + Inches(1.15),
             col_w - Inches(0.6), Inches(0.4),
             "writes JSON configs or raw numbers",
             size=12, italic=True, color=TEXT_MUTED)

    # JSON snippet with issues highlighted
    jy = y + Inches(1.70)
    jh = Inches(2.2)
    add_rect(s, left_x + Inches(0.3), jy, col_w - Inches(0.6), jh,
             fill=CODE_BG, line=None, rounded=True)
    json_bad = [
        ('{ "columns": [', CODE_FG),
        ('    { "name": "sales",', CODE_KEY),
        ('      "values": [42, 73, 61, 88]  },', CODE_STR),
        ('    { "name": "profit_margin",', CODE_KEY),
        ('      "values": [0.12, 0.21, 0.18,', CODE_STR),
        ('                 1.34]  ← ', CODE_STR),
        ('      }', CODE_FG),
        ('  ], "correlation": 0.9 }', CODE_FG),
    ]
    ly = jy + Inches(0.15)
    for text, col in json_bad:
        add_text(s, left_x + Inches(0.5), ly, col_w - Inches(0.9), Inches(0.22),
                 text, size=11, color=col, font=MONO, line_spacing=1.0)
        ly = ly + Inches(0.24)

    # problems list
    py = jy + jh + Inches(0.25)
    probs = [
        "✗  Numeric impossibilities slip through",
        "✗  Correlations infeasible to enforce",
        "✗  Only syntactic validation possible",
    ]
    for p in probs:
        add_text(s, left_x + Inches(0.4), py, col_w - Inches(0.8), Inches(0.28),
                 p, size=11, color=ACCENT, bold=True)
        py = py + Inches(0.30)

    # big arrow in the middle (between the two cards)
    arr_x = Inches(6.45)
    add_arrow_right(s, arr_x, y + h/2 - Inches(0.25), Inches(0.45), Inches(0.5),
                    fill=NAVY)

    # RIGHT — new way (featured)
    add_rect(s, right_x, y, col_w, h, fill=CARD_BG, line=NAVY, line_w=1.5, rounded=True)
    add_rect(s, right_x, y, col_w, Inches(0.12), fill=ACCENT, line=None)
    add_text(s, right_x + Inches(0.3), y + Inches(0.3),
             col_w - Inches(0.6), Inches(0.35),
             "AFTER", size=11, bold=True, color=ACCENT)
    add_text(s, right_x + Inches(0.3), y + Inches(0.6),
             col_w - Inches(0.6), Inches(0.6),
             "LLM-as-Data-Programmer", size=20, bold=True, color=NAVY)
    add_text(s, right_x + Inches(0.3), y + Inches(1.15),
             col_w - Inches(0.6), Inches(0.4),
             "writes Python calling a type-safe SDK",
             size=12, italic=True, color=TEXT_MID)

    cy = y + Inches(1.70)
    ch = Inches(2.2)
    add_rect(s, right_x + Inches(0.3), cy, col_w - Inches(0.6), ch,
             fill=CODE_BG, line=None, rounded=True)
    code_good = [
        ('sim.add_category("hospital",', CODE_FG, False),
        ('    values=["Xiehe","Huashan",...],', CODE_FG, False),
        ('    group="entity")', CODE_FG, False),
        ('', CODE_FG, False),
        ('sim.add_measure("wait_minutes",', CODE_FG, False),
        ('    family="lognormal",', CODE_FG, False),
        ('    param_model={...})', CODE_FG, False),
        ('', CODE_FG, False),
        ('sim.declare_orthogonal(', CODE_KEY, True),
        ('    "entity", "patient")', CODE_KEY, True),
    ]
    ly = cy + Inches(0.15)
    for text, col, bold in code_good:
        add_text(s, right_x + Inches(0.5), ly, col_w - Inches(0.9), Inches(0.22),
                 text, size=11, color=col, font=MONO, bold=bold, line_spacing=1.0)
        ly = ly + Inches(0.19)

    # benefits list
    by = cy + ch + Inches(0.25)
    bens = [
        "✓  Semantic errors caught at execution",
        "✓  Typed exceptions → targeted LLM repair",
        "✓  3-layer validator runs in milliseconds",
    ]
    for b in bens:
        add_text(s, right_x + Inches(0.4), by, col_w - Inches(0.8), Inches(0.28),
                 b, size=11, color=TEAL, bold=True)
        by = by + Inches(0.30)

    add_footer(s, 7)


# ─────────────────────────── Slide 8 — SDK Cheatsheet ───────────────────────────

def slide_sdk(prs):
    s = blank_slide(prs)
    add_title_bar(s, "The FactTableSimulator SDK",
                  "Minimal, strongly-typed API — everything an LLM needs, nothing it doesn't.")

    # LEFT — API table
    left_x = Inches(0.55); col_w = Inches(6.3)
    y = Inches(1.75); h = Inches(5.1)
    add_rect(s, left_x, y, col_w, h, fill=CARD_BG, line=BORDER, line_w=0.75, rounded=True)
    add_text(s, left_x + Inches(0.25), y + Inches(0.20),
             col_w - Inches(0.5), Inches(0.4),
             "SDK methods", size=15, bold=True, color=NAVY)

    # Step 1 group
    sy = y + Inches(0.70)
    add_text(s, left_x + Inches(0.25), sy, col_w - Inches(0.5), Inches(0.22),
             "STEP 1  ·  Column declarations",
             size=10, bold=True, color=ACCENT)
    api_step1 = [
        ("add_category",            "cat col  •  group  •  hierarchy via parent"),
        ("add_temporal",            "time col  •  auto-derive DOW / month"),
        ("add_measure",             "stochastic root  •  dist + effects"),
        ("add_measure_structural",  "formula of other measures + noise"),
    ]
    row_h = Inches(0.37)
    row_step = Inches(0.42)
    ly = sy + Inches(0.28)
    for name, desc in api_step1:
        add_rect(s, left_x + Inches(0.25), ly, col_w - Inches(0.5), row_h,
                 fill=SOFT_BG, line=None, rounded=True)
        add_text(s, left_x + Inches(0.40), ly, Inches(2.4), row_h,
                 name, size=11, bold=True, color=NAVY, font=MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left_x + Inches(2.85), ly, col_w - Inches(3.2), row_h,
                 desc, size=10.5, color=TEXT_MID, anchor=MSO_ANCHOR.MIDDLE)
        ly = ly + row_step

    # Step 2
    sy = ly + Inches(0.10)
    add_text(s, left_x + Inches(0.25), sy, col_w - Inches(0.5), Inches(0.22),
             "STEP 2  ·  Relationships & patterns",
             size=10, bold=True, color=ACCENT)
    api_step2 = [
        ("declare_orthogonal",    "group ⊥ group  →  propagates to pairs"),
        ("add_group_dependency",  "root-level DAG across groups"),
        ("inject_pattern",        "outlier / trend_break / reversal / ..."),
        ("set_realism",           "missing / dirty / censoring  (optional)"),
    ]
    ly = sy + Inches(0.28)
    for name, desc in api_step2:
        add_rect(s, left_x + Inches(0.25), ly, col_w - Inches(0.5), row_h,
                 fill=SOFT_BG, line=None, rounded=True)
        add_text(s, left_x + Inches(0.40), ly, Inches(2.4), row_h,
                 name, size=11, bold=True, color=NAVY, font=MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left_x + Inches(2.85), ly, col_w - Inches(3.2), row_h,
                 desc, size=10.5, color=TEXT_MID, anchor=MSO_ANCHOR.MIDDLE)
        ly = ly + row_step

    # RIGHT — code example
    right_x = Inches(7.05); col2_w = Inches(6.0)
    add_rect(s, right_x, y, col2_w, h, fill=CODE_BG, line=None, rounded=True)
    add_text(s, right_x + Inches(0.3), y + Inches(0.2),
             col2_w - Inches(0.6), Inches(0.3),
             "emergency_records.py", size=10, bold=True, color=CODE_KEY, font=MONO)

    code_lines = [
        ('sim = FactTableSimulator(target_rows=500)', CODE_FG),
        ('', CODE_FG),
        ('# Step 1 — declare columns', CODE_KEY),
        ('sim.add_category("hospital",', CODE_FG),
        ('   values=["Xiehe","Huashan","Ruijin",...],', CODE_FG),
        ('   group="entity")', CODE_FG),
        ('', CODE_FG),
        ('sim.add_category("severity",', CODE_FG),
        ('   values=["Mild","Moderate","Severe"],', CODE_FG),
        ('   group="patient")', CODE_FG),
        ('', CODE_FG),
        ('sim.add_measure("wait_minutes",', CODE_FG),
        ('   family="lognormal",', CODE_FG),
        ('   param_model={', CODE_FG),
        ('     "mu": {"intercept": 2.8,', CODE_FG),
        ('       "effects": {"severity": {...}}}})', CODE_FG),
        ('', CODE_FG),
        ('sim.add_measure_structural("cost",', CODE_FG),
        ('   formula="wait_minutes*12 + surcharge",', CODE_FG),
        ('   noise={"family":"gaussian","sigma":30})', CODE_FG),
        ('', CODE_FG),
        ('# Step 2 — relationships & patterns', CODE_KEY),
        ('sim.declare_orthogonal("entity","patient")', CODE_STR),
        ('sim.inject_pattern("outlier_entity", ...)', CODE_STR),
        ('', CODE_FG),
        ('df, meta = sim.generate()', CODE_STR),
    ]
    ly = y + Inches(0.55)
    for text, col in code_lines:
        add_text(s, right_x + Inches(0.35), ly, col2_w - Inches(0.7), Inches(0.2),
                 text, size=10.5, color=col, font=MONO, line_spacing=1.0)
        ly = ly + Inches(0.168)

    add_footer(s, 8)


# ─────────────────────────── Slide 9 — Dimension Groups ───────────────────────────

def slide_groups(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Dimension groups & orthogonality",
                  "One abstraction unifies categorical, temporal, cross-group semantics.")

    # Three group boxes
    y = Inches(1.85)
    gh = Inches(3.4)
    gw = Inches(3.7)
    gap = Inches(0.25)
    xs = [Inches(0.55) + i*(gw + gap) for i in range(3)]

    groups = [
        {
            "name": "entity",
            "color": NAVY,
            "items": [("hospital",   "root, 5 vals"),
                      ("department", "child, 4 vals"),
                      ("ward",       "grand-child")],
        },
        {
            "name": "patient",
            "color": TEAL,
            "items": [("severity",     "root, 3 vals"),
                      ("acuity_level", "child, 4 vals")],
        },
        {
            "name": "time",
            "color": ACCENT,
            "items": [("visit_date",  "root (temporal)"),
                      ("day_of_week", "derived"),
                      ("month",       "derived")],
        },
    ]

    for x, g in zip(xs, groups):
        add_rect(s, x, y, gw, gh, fill=CARD_BG, line=BORDER, line_w=0.75, rounded=True)
        # header
        add_rect(s, x, y, gw, Inches(0.55), fill=g["color"], line=None)
        add_text(s, x + Inches(0.25), y + Inches(0.12),
                 gw - Inches(0.5), Inches(0.35),
                 f'group  "{g["name"]}"', size=14, bold=True, color=WHITE, font=MONO)

        # items as nested boxes (hierarchy)
        ix = x + Inches(0.35); iy = y + Inches(0.85)
        for depth, (name, note) in enumerate(g["items"]):
            offset = Inches(0.28 * depth)
            # vertical connector line
            if depth > 0:
                add_line(s, ix + offset - Inches(0.15), iy - Inches(0.10),
                         ix + offset - Inches(0.15), iy + Inches(0.22),
                         color=BORDER, width_pt=1.0)
                add_line(s, ix + offset - Inches(0.15), iy + Inches(0.22),
                         ix + offset, iy + Inches(0.22),
                         color=BORDER, width_pt=1.0)
            bw = gw - offset - Inches(0.5)
            add_rect(s, ix + offset, iy, bw, Inches(0.55),
                     fill=SOFT_BG, line=BORDER, line_w=0.5, rounded=True)
            add_text(s, ix + offset + Inches(0.15), iy + Inches(0.04),
                     bw - Inches(0.25), Inches(0.3),
                     name, size=12, bold=True, color=NAVY, font=MONO)
            add_text(s, ix + offset + Inches(0.15), iy + Inches(0.28),
                     bw - Inches(0.25), Inches(0.25),
                     note, size=9, italic=True, color=TEXT_MUTED)
            iy = iy + Inches(0.68)

    # orthogonality line between group 1 and group 2
    g1_cx = xs[0] + gw / 2
    g2_cx = xs[1] + gw / 2
    y_orth = Inches(5.50)
    add_line(s, g1_cx, y + gh, g1_cx, y_orth, color=ACCENT, width_pt=1.5, dashed=True)
    add_line(s, g2_cx, y + gh, g2_cx, y_orth, color=ACCENT, width_pt=1.5, dashed=True)
    add_line(s, g1_cx, y_orth, g2_cx, y_orth, color=ACCENT, width_pt=1.5, dashed=True)
    # orth label bubble
    bubble_w = Inches(2.8)
    add_rect(s, (g1_cx + g2_cx)/2 - bubble_w/2, y_orth - Inches(0.24),
             bubble_w, Inches(0.48), fill=CARD_BG, line=ACCENT, line_w=1.5, rounded=True)
    add_text(s, (g1_cx + g2_cx)/2 - bubble_w/2, y_orth - Inches(0.24),
             bubble_w, Inches(0.48),
             "entity  ⊥  patient   (declare_orthogonal)",
             size=12, bold=True, color=ACCENT, font=MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # bottom callouts
    c_y = Inches(6.25)
    add_text(s, Inches(0.55), c_y, Inches(12.2), Inches(0.3),
             "Declare once at group level  →  propagates to all cross-group pairs  →  validated by χ² test",
             size=12, italic=True, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.55), c_y + Inches(0.35), Inches(12.2), Inches(0.3),
             "→  Enables Orthogonal-Slice dashboards & independence-based multi-chart QA",
             size=11, color=TEXT_MUTED, italic=True, align=PP_ALIGN.CENTER)

    add_footer(s, 9)


# ─────────────────────────── Slide 10 — DAG-ordered generation ───────────────────────────

def slide_dag(prs):
    s = blank_slide(prs)
    add_title_bar(s, "DAG-ordered event-level generation",
                  "Every row is one atomic event.  No cross-product, no repetition.")

    # Four layers
    layer_x = Inches(0.55)
    layer_w = Inches(12.22)
    y0 = Inches(1.75)
    layer_h = Inches(1.02)
    gap = Inches(0.12)

    layers = [
        {
            "tag": "LAYER 0",
            "sub": "Independent roots",
            "color": NAVY,
            "cols": [("hospital",    "Cat"),
                     ("severity",    "Cat"),
                     ("visit_date",  "Temp")],
        },
        {
            "tag": "LAYER 1",
            "sub": "Dependent non-measures",
            "color": TEAL,
            "cols": [("department",  "Cat | hospital"),
                     ("payment",     "Cat | severity"),
                     ("day_of_week", "derive(visit_date)"),
                     ("month",       "derive(visit_date)")],
        },
        {
            "tag": "LAYER 2",
            "sub": "Stochastic measure",
            "color": ACCENT,
            "cols": [("wait_minutes", "LogNormal(μ(sev, hosp), σ(sev))")],
        },
        {
            "tag": "LAYER 3",
            "sub": "Structural measures",
            "color": RGBColor(0x7B, 0x3F, 0xA0),
            "cols": [("cost",         "12·wait + surcharge(sev) + ε"),
                     ("satisfaction", "9 − 0.04·wait + adj(sev) + ε")],
        },
    ]

    # dim col widths
    for i, L in enumerate(layers):
        ly = y0 + (layer_h + gap) * i
        add_rect(s, layer_x, ly, layer_w, layer_h, fill=CARD_BG,
                 line=BORDER, line_w=0.5, rounded=True)
        # left tag band
        add_rect(s, layer_x, ly, Inches(1.8), layer_h, fill=L["color"], line=None)
        add_text(s, layer_x + Inches(0.15), ly + Inches(0.15),
                 Inches(1.55), Inches(0.32),
                 L["tag"], size=11, bold=True, color=WHITE)
        add_text(s, layer_x + Inches(0.15), ly + Inches(0.45),
                 Inches(1.55), Inches(0.45),
                 L["sub"], size=10, color=WHITE, italic=True, line_spacing=1.1)

        # columns as chips
        ccx = layer_x + Inches(2.0)
        cols = L["cols"]
        n = len(cols)
        avail_w = layer_w - Inches(2.15) - Inches(0.15)
        cw = (avail_w - Inches(0.15) * (n - 1)) / n if n > 0 else avail_w
        for i, (name, rhs) in enumerate(cols):
            cx = ccx + (cw + Inches(0.15)) * i
            add_rect(s, cx, ly + Inches(0.14), cw, layer_h - Inches(0.28),
                     fill=SOFT_BG, line=BORDER, line_w=0.5, rounded=True)
            add_text(s, cx + Inches(0.15), ly + Inches(0.18),
                     cw - Inches(0.25), Inches(0.3),
                     name, size=11.5, bold=True, color=NAVY, font=MONO)
            add_text(s, cx + Inches(0.15), ly + Inches(0.45),
                     cw - Inches(0.25), Inches(0.4),
                     rhs, size=10, italic=True, color=TEXT_MID, font=MONO,
                     line_spacing=1.15)

    # down arrows between layers
    for i in range(3):
        ax = Inches(0.55) + Inches(0.9) - Inches(0.08)
        ay = y0 + (layer_h + gap) * i + layer_h - Inches(0.02)
        add_arrow_right(
            s,
            Inches(0.55) + Inches(0.9) - Inches(0.10),
            y0 + (layer_h + gap) * i + layer_h - Inches(0.02),
            Inches(0.20), Inches(0.18), fill=TEXT_MID,
        )
        # replace with downward arrow using shape rotation
    # Use a plain down-arrow shape at center of tag band between layers
    # (we already added small right arrows; remove and add proper downward ones)

    # bottom callout: post-generation
    pg_y = Inches(6.4)
    add_rect(s, Inches(0.55), pg_y, Inches(12.22), Inches(0.55),
             fill=GRAY_BG, line=None, rounded=True)
    add_text(s, Inches(0.75), pg_y + Inches(0.1),
             Inches(12), Inches(0.35),
             "After DAG-ordered sampling:  γ  pattern injection  →  δ  realism (optional)  →  τ  post-process",
             size=12, bold=True, color=TEXT_DARK, font=MONO)

    add_footer(s, 10)


# ─────────────────────────── Slide 11 — Three-layer validation ───────────────────────────

def slide_validation(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Three-layer validation  +  code-execution feedback",
                  "Semantic errors, not just syntactic ones.  Auto-fix without LLM re-call.")

    # Top horizontal flow: LLM → Sandbox → Validator → Fix
    fy = Inches(1.85)
    fh = Inches(0.9)
    node_w = Inches(2.7)
    gap = Inches(0.3)
    total_w = node_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2

    nodes = [
        ("LLM", "writes Python\ncalling the SDK", NAVY),
        ("Sandbox", "executes\nbuild_fact_table()", TEAL),
        ("Validator", "L1 · L2 · L3\ndeterministic checks", ACCENT),
        ("Auto-fix", "parameter tweak\nno LLM re-call", RGBColor(0x7B, 0x3F, 0xA0)),
    ]
    for i, (title, body, col) in enumerate(nodes):
        nx = start_x + (node_w + gap) * i
        add_rect(s, nx, fy, node_w, fh, fill=CARD_BG, line=col, line_w=1.5, rounded=True)
        add_rect(s, nx, fy, node_w, Inches(0.24), fill=col, line=None)
        add_text(s, nx + Inches(0.15), fy + Inches(0.02),
                 node_w - Inches(0.3), Inches(0.25),
                 title, size=11, bold=True, color=WHITE)
        add_text(s, nx + Inches(0.15), fy + Inches(0.30),
                 node_w - Inches(0.3), Inches(0.55),
                 body, size=11, color=TEXT_DARK, line_spacing=1.2)
        if i < 3:
            ax = nx + node_w + Inches(0.02)
            add_arrow_right(s, ax, fy + fh/2 - Inches(0.10),
                            gap - Inches(0.04), Inches(0.2), fill=NAVY)

    # Feedback arc from Auto-fix back to Sandbox (below the flow)
    af_cx = start_x + (node_w + gap) * 3 + node_w / 2
    sb_cx = start_x + (node_w + gap) * 1 + node_w / 2
    arc_y = fy + fh + Inches(0.35)
    add_line(s, af_cx, fy + fh, af_cx, arc_y, color=ACCENT, width_pt=1.5, dashed=True)
    add_line(s, af_cx, arc_y, sb_cx, arc_y, color=ACCENT, width_pt=1.5, dashed=True)
    add_line(s, sb_cx, arc_y, sb_cx, fy + fh, color=ACCENT, width_pt=1.5, dashed=True, arrow_end=True)
    add_text(s, (sb_cx + af_cx)/2 - Inches(2), arc_y + Inches(0.03),
             Inches(4), Inches(0.3),
             "re-execute  (max 3 retries)",
             size=10, italic=True, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)

    # Failure feedback to LLM (further below, shown separately)
    llm_cx = start_x + node_w / 2
    val_cx = start_x + (node_w + gap) * 2 + node_w / 2
    arc_y2 = arc_y + Inches(0.40)
    add_line(s, val_cx, arc_y, val_cx, arc_y2, color=TEXT_MUTED, width_pt=1.2, dashed=True)
    add_line(s, val_cx, arc_y2, llm_cx, arc_y2, color=TEXT_MUTED, width_pt=1.2, dashed=True)
    add_line(s, llm_cx, arc_y2, llm_cx, fy + fh, color=TEXT_MUTED, width_pt=1.2, dashed=True, arrow_end=True)
    add_text(s, (llm_cx + val_cx)/2 - Inches(2.3), arc_y2 + Inches(0.03),
             Inches(4.6), Inches(0.3),
             "typed exception + traceback  →  LLM repair",
             size=10, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Three validation layer cards
    cy = Inches(3.95)
    ch = Inches(2.95)
    cw = Inches(4.0)
    cgap = Inches(0.11)
    cxs = [Inches(0.55) + i*(cw + cgap) for i in range(3)]
    layers = [
        {
            "label": "L1",
            "title": "Structural",
            "color": NAVY,
            "checks": [
                "Row count within 10%",
                "Categorical cardinality",
                "Finite measures",
                "χ²  for orthogonal pairs",
                "Measure DAG acyclicity",
            ],
        },
        {
            "label": "L2",
            "title": "Statistical",
            "color": TEAL,
            "checks": [
                "KS test per predictor cell",
                "Structural residual mean",
                "Structural residual σ",
                "Conditional transitions",
                "Correlation ± 0.15",
            ],
        },
        {
            "label": "L3",
            "title": "Pattern",
            "color": ACCENT,
            "checks": [
                "Outlier z ≥ 2.0",
                "Trend-break magnitude",
                "Ranking reversal sign",
                "Dominance-shift check",
                "Seasonal anomaly",
            ],
        },
    ]
    for x, L in zip(cxs, layers):
        add_rect(s, x, cy, cw, ch, fill=CARD_BG, line=BORDER, line_w=0.75, rounded=True)
        # big tag
        add_rect(s, x + Inches(0.25), cy + Inches(0.22), Inches(0.75), Inches(0.75),
                 fill=L["color"], line=None, rounded=True)
        add_text(s, x + Inches(0.25), cy + Inches(0.22), Inches(0.75), Inches(0.75),
                 L["label"], size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.15), cy + Inches(0.32),
                 cw - Inches(1.35), Inches(0.4),
                 L["title"], size=17, bold=True, color=NAVY)
        # checks
        ly = cy + Inches(1.20)
        for chk in L["checks"]:
            # check glyph
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.35), ly + Inches(0.10),
                                     Inches(0.10), Inches(0.10))
            _set_fill(dot, L["color"]); _set_line(dot, invisible=True)
            add_text(s, x + Inches(0.55), ly, cw - Inches(0.8), Inches(0.3),
                     chk, size=11, color=TEXT_DARK)
            ly = ly + Inches(0.33)

    add_footer(s, 11)


# ─────────────────────────── Slide 12 — Table amortization ───────────────────────────

def slide_amortization(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Table amortization: one table → many tasks",
                  "Every derived chart shares the same ground-truth arithmetic by construction.")

    # Center: Master Table node (left)
    mx, my, mw, mh = Inches(0.7), Inches(2.7), Inches(3.4), Inches(2.3)
    add_rect(s, mx, my, mw, mh, fill=NAVY, line=None, rounded=True)
    add_rect(s, mx, my, mw, Inches(0.5), fill=NAVY_DARK, line=None)
    add_text(s, mx + Inches(0.3), my + Inches(0.10), mw - Inches(0.6), Inches(0.3),
             "MASTER FACT TABLE", size=12, bold=True, color=WHITE)
    # rows illustration
    for i in range(6):
        add_rect(s, mx + Inches(0.3), my + Inches(0.75) + Inches(0.22)*i,
                 mw - Inches(0.6), Inches(0.15),
                 fill=WHITE if i % 2 else CHIP_BG, line=None)
    # caption
    add_text(s, mx, my + mh + Inches(0.1), mw, Inches(0.3),
             "atomic events  ·  500–3000 rows",
             size=11, italic=True, color=TEXT_MID,
             align=PP_ALIGN.CENTER)

    # Formula
    fy = Inches(1.80)
    add_rect(s, Inches(0.7), fy, Inches(3.4), Inches(0.75),
             fill=GRAY_BG, line=BORDER, line_w=0.5, rounded=True)
    add_text(s, Inches(0.7), fy + Inches(0.10), Inches(3.4), Inches(0.25),
             "VIEW  =", size=11, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), fy + Inches(0.36), Inches(3.4), Inches(0.35),
             "σ ∘ γ ∘ π  (M)",
             size=20, bold=True, color=NAVY, font=MONO,
             align=PP_ALIGN.CENTER)

    # Arrows fanning out from master
    hub_x = mx + mw
    hub_y = my + mh / 2

    # Right panel: 6 chart families
    right_x = Inches(5.0)
    fam_w = Inches(3.95)
    fam_h = Inches(0.72)
    fam_gap = Inches(0.12)
    right_y = Inches(1.75)
    families = [
        ("Comparison",  "bar  ·  grouped_bar",                          NAVY),
        ("Trend",       "line  ·  area",                                TEAL),
        ("Distribution","histogram  ·  box  ·  violin",                 ACCENT),
        ("Composition", "pie  ·  donut  ·  stacked_bar  ·  treemap",    RGBColor(0x7B, 0x3F, 0xA0)),
        ("Relationship","scatter  ·  bubble  ·  heatmap  ·  radar",     RGBColor(0x1B, 0x6F, 0x4E)),
        ("Flow",        "waterfall  ·  funnel",                         RGBColor(0xB5, 0x41, 0x48)),
    ]
    for i, (name, types, col) in enumerate(families):
        fy = right_y + (fam_h + fam_gap) * i
        add_rect(s, right_x, fy, fam_w, fam_h, fill=CARD_BG, line=BORDER, line_w=0.5, rounded=True)
        # color tab
        add_rect(s, right_x, fy, Inches(0.12), fam_h, fill=col, line=None)
        add_text(s, right_x + Inches(0.3), fy + Inches(0.08),
                 fam_w - Inches(0.5), Inches(0.3),
                 name, size=12, bold=True, color=NAVY)
        add_text(s, right_x + Inches(0.3), fy + Inches(0.38),
                 fam_w - Inches(0.45), Inches(0.3),
                 types, size=9.5, color=TEXT_MID, font=MONO)
        # fan arrow from master
        add_line(s, hub_x, hub_y, right_x - Inches(0.02),
                 fy + fam_h / 2, color=TEXT_MUTED, width_pt=0.75, arrow_end=True)

    # Right-rightmost: relationship types & count badge
    far_x = Inches(9.3)
    far_y = Inches(1.75)
    far_w = Inches(3.48)
    far_h = Inches(5.1)
    add_rect(s, far_x, far_y, far_w, far_h, fill=SOFT_BG, line=BORDER, line_w=0.75, rounded=True)
    add_text(s, far_x + Inches(0.25), far_y + Inches(0.15),
             far_w - Inches(0.5), Inches(0.4),
             "7 inter-chart relationships",
             size=14, bold=True, color=NAVY)
    rels = [
        "Drill-down",
        "Orthogonal Slice",
        "Comparative",
        "Dual-Metric",
        "Part-Whole",
        "Associative",
        "Causal Chain",
    ]
    ly = far_y + Inches(0.70)
    for r in rels:
        # bullet
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 far_x + Inches(0.35), ly + Inches(0.10),
                                 Inches(0.10), Inches(0.10))
        _set_fill(dot, ACCENT); _set_line(dot, invisible=True)
        add_text(s, far_x + Inches(0.55), ly, far_w - Inches(0.8), Inches(0.3),
                 r, size=12, color=TEXT_DARK)
        ly = ly + Inches(0.38)

    # big output number
    oy = ly + Inches(0.3)
    add_rect(s, far_x + Inches(0.35), oy, far_w - Inches(0.7), Inches(1.1),
             fill=NAVY, line=None, rounded=True)
    add_text(s, far_x + Inches(0.35), oy + Inches(0.05), far_w - Inches(0.7), Inches(0.5),
             "10 – 30+", size=30, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, far_x + Inches(0.35), oy + Inches(0.60), far_w - Inches(0.7), Inches(0.4),
             "coherent tasks per table",
             size=11, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # bottom quote
    add_text(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.3),
             "1 LLM call  →  1 script  →  1 Master Table  →  16 chart types × 7 relationships",
             size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_footer(s, 12)


# ─────────────────────────── Slide 13 — Operator algebra ───────────────────────────

def slide_algebra(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Operator algebra: a question is a typed pipeline",
                  "Two types  (V, S).  Four operator families.  Difficulty = # ops.")

    # Four operator family cards
    y = Inches(1.75)
    ch = Inches(2.0)
    cw = Inches(3.05)
    cgap = Inches(0.11)
    xs = [Inches(0.55) + i*(cw + cgap) for i in range(4)]
    fams = [
        {
            "tag": "SET  (V → V)",
            "color": NAVY,
            "ops": ["Filter", "Sort", "Limit", "GroupBy"],
        },
        {
            "tag": "SCALAR  (V → S)",
            "color": TEAL,
            "ops": ["Max", "Min", "Avg", "Sum", "Count", "ArgMax", "ArgMin", "ValueAt"],
        },
        {
            "tag": "COMBINATOR  (S,S → S)",
            "color": RGBColor(0x7B, 0x3F, 0xA0),
            "ops": ["Diff", "Ratio"],
        },
        {
            "tag": "BRIDGE  (S,V → V | V,V → S)",
            "color": ACCENT,
            "ops": ["EntityTransfer", "ValueTransfer", "TrendCompare", "RankCompare"],
        },
    ]
    for x, f in zip(xs, fams):
        add_rect(s, x, y, cw, ch, fill=CARD_BG, line=BORDER, line_w=0.75, rounded=True)
        add_rect(s, x, y, cw, Inches(0.45), fill=f["color"], line=None)
        add_text(s, x + Inches(0.15), y + Inches(0.08),
                 cw - Inches(0.3), Inches(0.35),
                 f["tag"], size=10.5, bold=True, color=WHITE, font=MONO)
        # op chips
        is_bridge = (f["tag"].startswith("BRIDGE"))
        chip_size = 9 if is_bridge else 10
        chip_h = 0.30 if is_bridge else 0.32
        wrap_step = Inches(0.36) if is_bridge else Inches(0.38)
        ox = x + Inches(0.18); oy = y + Inches(0.6)
        for op in f["ops"]:
            if is_bridge:
                # render each bridge op on its own row, full card width
                bw = cw - Inches(0.36)
                add_rect(s, ox, oy, bw, Inches(chip_h),
                         fill=SOFT_BG, line=None, rounded=True)
                tf_shape = s.shapes[-1]
                tf_shape.adjustments[0] = 0.5
                add_text(s, ox, oy, bw, Inches(chip_h),
                         op, size=chip_size, bold=False,
                         color=NAVY, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
                oy = oy + wrap_step
            else:
                shp, ow = add_chip(s, ox, oy, op, size=chip_size,
                                   fill=SOFT_BG, color=NAVY, height=chip_h)
                ox = ox + ow + Inches(0.06)
                if ox + Inches(0.7) > x + cw - Inches(0.15):
                    ox = x + Inches(0.18); oy = oy + wrap_step

    # Pipeline example
    py = Inches(4.0)
    add_text(s, Inches(0.55), py, Inches(12.2), Inches(0.35),
             "Example pipeline  (3 ops, MEDIUM)",
             size=12, bold=True, color=ACCENT)

    # Pipeline boxes + arrows
    pipe_y = Inches(4.5)
    pipe_h = Inches(0.8)
    # order: V  →  Sort  →  V  →  Limit(3)  →  V  →  Avg  →  S
    items = [
        ("V",        NAVY,   True),
        ("Sort",     WHITE,  False),
        ("V",        NAVY,   True),
        ("Limit(3)", WHITE,  False),
        ("V",        NAVY,   True),
        ("Avg",      WHITE,  False),
        ("S",        ACCENT, True),
    ]
    total_items = len(items)
    # widths: types narrower, ops wider
    ws = [Inches(0.7) if is_type else Inches(1.4) for (_, _, is_type) in items]
    arr_w = Inches(0.45)
    tot = sum(w.emu for w in ws) + arr_w.emu * (total_items - 1)
    start_x = (SLIDE_W.emu - tot) // 2
    cur_x = Emu(start_x)
    for i, ((lbl, col, is_type), w) in enumerate(zip(items, ws)):
        if is_type:
            add_rect(s, cur_x, pipe_y, w, pipe_h, fill=col, line=None, rounded=True)
            add_text(s, cur_x, pipe_y, w, pipe_h,
                     lbl, size=20, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        else:
            add_rect(s, cur_x, pipe_y, w, pipe_h, fill=col,
                     line=NAVY, line_w=1.2, rounded=True)
            add_text(s, cur_x, pipe_y, w, pipe_h,
                     lbl, size=14, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        cur_x = Emu(cur_x.emu + w.emu)
        if i < total_items - 1:
            add_arrow_right(s, cur_x, pipe_y + pipe_h/2 - Inches(0.12),
                            arr_w, Inches(0.24), fill=TEXT_MID)
            cur_x = Emu(cur_x.emu + arr_w.emu)

    # Question box
    qy = Inches(5.85)
    add_rect(s, Inches(1.5), qy, Inches(10.3), Inches(0.95),
             fill=SOFT_BG, line=BORDER, line_w=0.75, rounded=True)
    add_text(s, Inches(1.7), qy + Inches(0.10), Inches(10), Inches(0.35),
             "QUESTION", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(1.7), qy + Inches(0.38), Inches(10), Inches(0.5),
             "“What is the average wait time of the top-3 hospitals?”",
             size=16, italic=True, color=NAVY, bold=True)

    add_footer(s, 13)


# ─────────────────────────── Slide 14 — Multi-chart bridges ───────────────────────────

def slide_bridges(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Cross-chart reasoning via bridge operators",
                  "A multi-chart question chains two pipelines through a typed bridge.")

    # LEFT — relationships + bridges table
    lx, ly, lw, lh = Inches(0.55), Inches(1.75), Inches(5.4), Inches(5.1)
    add_rect(s, lx, ly, lw, lh, fill=CARD_BG, line=BORDER, line_w=0.75, rounded=True)
    add_text(s, lx + Inches(0.25), ly + Inches(0.15),
             lw - Inches(0.5), Inches(0.4),
             "7 relationships  →  valid bridges",
             size=13, bold=True, color=NAVY)

    rows = [
        ("Drill-down",        "EntityTransfer · ValueTransfer"),
        ("Orthogonal Slice",  "EntityTransfer · RankCompare"),
        ("Comparative",       "ValueTransfer · TrendCompare"),
        ("Dual-Metric",       "EntityTransfer · RankCompare · ValueTransfer"),
        ("Part-Whole",        "EntityTransfer · ValueTransfer"),
        ("Associative",       "RankCompare · EntityTransfer"),
        ("Causal Chain",      "EntityTransfer · ValueTransfer"),
    ]
    rh = Inches(0.52)
    ry = ly + Inches(0.70)
    for i, (rel, br) in enumerate(rows):
        bgc = SOFT_BG if i % 2 else WHITE
        add_rect(s, lx + Inches(0.25), ry, lw - Inches(0.5), rh,
                 fill=bgc, line=BORDER, line_w=0.4, rounded=False)
        add_text(s, lx + Inches(0.4), ry, Inches(1.9), rh,
                 rel, size=11.5, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lx + Inches(2.3), ry, lw - Inches(2.6), rh,
                 br, size=10.5, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        ry = ry + rh

    # RIGHT — concrete example
    rx = Inches(6.20); rw = Inches(6.58); ry0 = Inches(1.75); rh0 = Inches(5.1)
    add_rect(s, rx, ry0, rw, rh0, fill=CARD_BG, line=NAVY, line_w=1.2, rounded=True)
    add_rect(s, rx, ry0, rw, Inches(0.12), fill=ACCENT, line=None)
    add_text(s, rx + Inches(0.25), ry0 + Inches(0.2),
             rw - Inches(0.5), Inches(0.35),
             "EXAMPLE  ·  Dual-Metric with EntityTransfer",
             size=11, bold=True, color=ACCENT)

    # Two mini chart cards
    chart_h = Inches(1.35)
    chart_w = Inches(2.6)
    ca_x = rx + Inches(0.4); cb_x = rx + rw - Inches(0.4) - chart_w
    cay = ry0 + Inches(0.8)

    # Chart A: bar hospital × wait
    add_rect(s, ca_x, cay, chart_w, chart_h, fill=SOFT_BG, line=BORDER, line_w=0.75, rounded=True)
    add_text(s, ca_x + Inches(0.15), cay + Inches(0.1),
             chart_w - Inches(0.3), Inches(0.3),
             "Chart A: hospital × wait", size=10, bold=True, color=NAVY)
    bars_a = [0.55, 0.85, 0.45, 0.65]
    for i, hgt in enumerate(bars_a):
        bx = ca_x + Inches(0.25 + i*0.5)
        by = cay + chart_h - Inches(0.18) - Inches(hgt)
        col = ACCENT if hgt == max(bars_a) else NAVY
        add_rect(s, bx, by, Inches(0.35), Inches(hgt), fill=col, line=None)
    # labels
    add_text(s, ca_x + Inches(0.25 + 1*0.5) - Inches(0.25), cay + chart_h - Inches(0.17),
             Inches(0.8), Inches(0.18),
             "Xiehe", size=8, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Chart B: bar hospital × cost
    add_rect(s, cb_x, cay, chart_w, chart_h, fill=SOFT_BG, line=BORDER, line_w=0.75, rounded=True)
    add_text(s, cb_x + Inches(0.15), cay + Inches(0.1),
             chart_w - Inches(0.3), Inches(0.3),
             "Chart B: hospital × cost", size=10, bold=True, color=NAVY)
    bars_b = [0.50, 0.70, 0.40, 0.55]
    for i, hgt in enumerate(bars_b):
        bx = cb_x + Inches(0.25 + i*0.5)
        by = cay + chart_h - Inches(0.18) - Inches(hgt)
        col = ACCENT if i == 1 else TEAL
        add_rect(s, bx, by, Inches(0.35), Inches(hgt), fill=col, line=None)
    add_text(s, cb_x + Inches(0.25 + 1*0.5) - Inches(0.25), cay + chart_h - Inches(0.17),
             Inches(0.8), Inches(0.18),
             "Xiehe", size=8, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Pipeline chain — rendered compactly inside the card
    py2 = cay + chart_h + Inches(0.30)
    chain_items = [
        ("V_a",             NAVY,   True),
        ("ArgMax",          WHITE,  False),
        ("'Xiehe'",         ACCENT, True),
        ("Bridge",          WHITE,  False),
        ("V_b",             NAVY,   True),
        ("ValueAt",         WHITE,  False),
        ("6200",            ACCENT, True),
    ]
    # compact widths based on text length
    ws = [Inches(max(0.48, 0.095 * len(lbl) + 0.28)) for (lbl, _, _) in chain_items]
    arr_w = Inches(0.18)
    box_h = Inches(0.48)
    tot = sum(w.emu for w in ws) + arr_w.emu * (len(chain_items) - 1)
    # center within the example card with small safety margin
    avail = rw.emu - Inches(0.5).emu
    if tot > avail:
        # uniformly scale down widths
        scale = avail / tot
        ws = [Emu(int(w.emu * scale)) for w in ws]
        arr_w = Emu(int(arr_w.emu * scale))
        tot = sum(w.emu for w in ws) + arr_w.emu * (len(chain_items) - 1)
    cx = Emu(rx.emu + (rw.emu - tot) // 2)
    for i, ((lbl, col, is_t), w) in enumerate(zip(chain_items, ws)):
        if is_t:
            add_rect(s, cx, py2, w, box_h,
                     fill=col, line=None, rounded=True)
            add_text(s, cx, py2, w, box_h,
                     lbl, size=10.5, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        else:
            add_rect(s, cx, py2, w, box_h,
                     fill=WHITE, line=NAVY, line_w=1.0, rounded=True)
            add_text(s, cx, py2, w, box_h,
                     lbl, size=9.5, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        cx = Emu(cx.emu + w.emu)
        if i < len(chain_items) - 1:
            add_arrow_right(s, cx, py2 + Inches(0.13),
                            arr_w, Inches(0.22), fill=TEXT_MID)
            cx = Emu(cx.emu + arr_w.emu)

    # Question box
    qy = py2 + Inches(0.75)
    add_rect(s, rx + Inches(0.35), qy, rw - Inches(0.7), Inches(0.85),
             fill=SOFT_BG, line=BORDER, line_w=0.5, rounded=True)
    add_text(s, rx + Inches(0.5), qy + Inches(0.08), rw - Inches(1), Inches(0.3),
             "QUESTION  ·  difficulty = 3 ops", size=10, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.5), qy + Inches(0.35), rw - Inches(1), Inches(0.45),
             "“The hospital with the longest wait — what is its cost?”",
             size=13, italic=True, color=NAVY, bold=True)

    add_footer(s, 14)


# ─────────────────────────── Slide 15 — Summary ───────────────────────────

def slide_summary(prs):
    s = blank_slide(prs)
    add_title_bar(s, "Summary",
                  "One formal substrate — programmable, amortized, deterministic.")

    # Four layer cards in a row
    y = Inches(1.85)
    h = Inches(4.2)
    w = Inches(3.0)
    gap = Inches(0.12)
    xs = [Inches(0.55) + i*(w + gap) for i in range(4)]
    layers = [
        {
            "tag": "LAYER 1",
            "color": NAVY,
            "title": "Operator Algebra",
            "idea": "Formalizes every chart QA as a typed pipeline.",
            "mech": "Set · Scalar · Combinator · Bridge  (16 ops, 2 types)",
        },
        {
            "tag": "LAYER 2",
            "color": TEAL,
            "title": "Agentic Simulator",
            "idea": "LLM writes executable SDK code — not JSON, not raw numbers.",
            "mech": "Type-safe API · DAG measures · Orthogonality · Pattern injection",
        },
        {
            "tag": "LAYER 3",
            "color": ACCENT,
            "title": "Table Amortization",
            "idea": "One Master Table serves 10–30+ coherent tasks.",
            "mech": "16 chart types  ×  7 inter-chart relationships",
        },
        {
            "tag": "LAYER 4",
            "color": RGBColor(0x7B, 0x3F, 0xA0),
            "title": "Rule-Based QA",
            "idea": "Deterministic, template-free generation with pattern-seeded hard QA.",
            "mech": "Intra-view · Inter-view · Pattern-triggered · Difficulty = # ops",
        },
    ]

    for x, L in zip(xs, layers):
        add_rect(s, x, y, w, h, fill=CARD_BG, line=BORDER, line_w=0.75, rounded=True)
        add_rect(s, x, y, w, Inches(0.10), fill=L["color"], line=None)
        add_text(s, x + Inches(0.25), y + Inches(0.25),
                 w - Inches(0.5), Inches(0.3),
                 L["tag"], size=10, bold=True, color=L["color"])
        add_text(s, x + Inches(0.25), y + Inches(0.55),
                 w - Inches(0.5), Inches(0.7),
                 L["title"], size=18, bold=True, color=NAVY, line_spacing=1.1)
        # divider
        add_line(s, x + Inches(0.25), y + Inches(1.45),
                 x + w - Inches(0.25), y + Inches(1.45),
                 color=BORDER)
        # idea
        add_text(s, x + Inches(0.25), y + Inches(1.60),
                 w - Inches(0.5), Inches(1.3),
                 L["idea"], size=12, color=TEXT_DARK, line_spacing=1.35)
        # mechanism label
        add_text(s, x + Inches(0.25), y + Inches(2.95),
                 w - Inches(0.5), Inches(0.3),
                 "MECHANISM", size=9, bold=True, color=TEXT_MUTED)
        add_text(s, x + Inches(0.25), y + Inches(3.25),
                 w - Inches(0.5), Inches(0.9),
                 L["mech"], size=10.5, italic=True, color=TEXT_MID,
                 line_spacing=1.3, font=MONO)

    # Bottom banner
    by = Inches(6.3)
    add_rect(s, Inches(0.55), by, Inches(12.22), Inches(0.75),
             fill=NAVY, line=None, rounded=True)
    add_text(s, Inches(0.55), by + Inches(0.05), Inches(12.22), Inches(0.35),
             "CORE GUARANTEE",
             size=10, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.55), by + Inches(0.33), Inches(12.22), Inches(0.45),
             "Every chart view shares the same ground-truth arithmetic by construction.",
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, italic=True)

    add_footer(s, 15)


# ─────────────────────────── Main ───────────────────────────

def main(out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_motivation(prs)
    slide_prior_work(prs)
    slide_contributions(prs)
    slide_pipeline(prs)
    slide_phase_01(prs)
    slide_paradigm(prs)
    slide_sdk(prs)
    slide_groups(prs)
    slide_dag(prs)
    slide_validation(prs)
    slide_amortization(prs)
    slide_algebra(prs)
    slide_bridges(prs)
    slide_summary(prs)

    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "chartagent_proposal.pptx"
    main(out)
